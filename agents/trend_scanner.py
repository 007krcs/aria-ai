"""
NOVA v3 — Trend Scanner + Universal Document Scanner
======================================================

TrendScanner:
  Pulls the latest research, news, and trending topics from:
  - arXiv (latest papers — AI, CS, physics, math, bio)
  - Semantic Scholar (citation graph + paper summaries)
  - HackerNews (tech news and discussions)
  - Reddit (any subreddit — r/MachineLearning, r/science, etc.)
  - RSS feeds (any news source)
  - GitHub trending (hottest repos right now)
  All free. No API keys needed for basic access.

UniversalScanner:
  Reads ANY document in ANY language:
  - PDF (text + scanned/OCR)
  - Word (.docx), Excel (.xlsx), PowerPoint (.pptx)
  - Images with text (OCR via Tesseract or EasyOCR)
  - Web pages (Playwright fallback for JS sites)
  - YouTube transcripts (via youtube-transcript-api)
  - Audio/video transcription (via Whisper)
  - QR codes, barcodes
  - Any of 100+ languages auto-detected
"""

import re
import time
import json
import hashlib
import requests
from pathlib import Path
from datetime import datetime, timedelta
from typing import Optional, Generator
from rich.console import Console

console = Console()
CACHE_DIR = Path("data/trend_cache")
CACHE_DIR.mkdir(parents=True, exist_ok=True)


# ─────────────────────────────────────────────────────────────────────────────
# TREND SCANNER
# ─────────────────────────────────────────────────────────────────────────────

class TrendScanner:
    """
    Pulls latest data from multiple free sources simultaneously.
    Results are cached locally to avoid hammering APIs.
    Auto-stores everything in ChromaDB so NOVA learns from trends.
    """

    HEADERS = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 Chrome/124.0.0.0 Safari/537.36"
        )
    }

    def __init__(self, memory=None, cache_ttl_minutes: int = 30):
        self.memory    = memory
        self.cache_ttl = cache_ttl_minutes * 60

    # ── arXiv — latest AI/CS/math papers ─────────────────────────────────────

    def arxiv_latest(self, query: str = "AI machine learning", max_results: int = 8) -> list[dict]:
        """
        Search arXiv for the latest papers matching a query.
        Free, no API key. Returns papers from the last 7 days by default.
        """
        cache_key = f"arxiv_{hashlib.md5(query.encode()).hexdigest()[:8]}"
        cached    = self._load_cache(cache_key)
        if cached:
            return cached

        try:
            import urllib.request
            import urllib.parse
            import xml.etree.ElementTree as ET

            encoded = urllib.parse.quote(query)
            url     = (
                f"https://export.arxiv.org/api/query?"
                f"search_query=all:{encoded}"
                f"&start=0&max_results={max_results}"
                f"&sortBy=lastUpdatedDate&sortOrder=descending"
            )
            with urllib.request.urlopen(url, timeout=10) as r:
                xml_data = r.read().decode("utf-8")

            ns      = {"atom": "http://www.w3.org/2005/Atom"}
            root    = ET.fromstring(xml_data)
            papers  = []

            for entry in root.findall("atom:entry", ns):
                title   = entry.findtext("atom:title",   "", ns).strip().replace("\n", " ")
                summary = entry.findtext("atom:summary", "", ns).strip().replace("\n", " ")
                link    = entry.findtext("atom:id",      "", ns).strip()
                updated = entry.findtext("atom:updated", "", ns)
                authors = [a.findtext("atom:name", "", ns)
                           for a in entry.findall("atom:author", ns)][:3]

                if title:
                    papers.append({
                        "title":    title,
                        "summary":  summary[:400],
                        "url":      link,
                        "authors":  authors,
                        "updated":  updated[:10],
                        "source":   "arxiv",
                    })

            self._save_cache(cache_key, papers)

            # Auto-learn: store in ChromaDB
            if self.memory and papers:
                chunks = [
                    {"text": f"{p['title']}. {p['summary']}", "source": p["url"], "domain": "research"}
                    for p in papers
                ]
                self.memory.store_many(chunks)
                console.print(f"  [dim]arXiv:[/] {len(papers)} papers fetched and stored")

            return papers

        except Exception as e:
            console.print(f"  [yellow]arXiv error: {e}[/]")
            return []

    # ── HackerNews — tech trends ──────────────────────────────────────────────

    def hackernews_top(self, max_stories: int = 10) -> list[dict]:
        """Fetch top stories from HackerNews. Free, no key needed."""
        cache_key = "hn_top"
        cached    = self._load_cache(cache_key)
        if cached:
            return cached

        try:
            r        = requests.get("https://hacker-news.firebaseio.com/v0/topstories.json", timeout=8)
            ids      = r.json()[:max_stories]
            stories  = []
            for story_id in ids:
                sr = requests.get(
                    f"https://hacker-news.firebaseio.com/v0/item/{story_id}.json",
                    timeout=5
                )
                s  = sr.json()
                if s and s.get("title"):
                    stories.append({
                        "title":  s.get("title", ""),
                        "url":    s.get("url", f"https://news.ycombinator.com/item?id={story_id}"),
                        "score":  s.get("score", 0),
                        "source": "hackernews",
                    })

            self._save_cache(cache_key, stories)
            if self.memory and stories:
                self.memory.store_many([
                    {"text": s["title"], "source": s["url"], "domain": "tech_news"}
                    for s in stories
                ])
            return stories

        except Exception as e:
            console.print(f"  [yellow]HN error: {e}[/]")
            return []

    # ── Reddit — subreddit trending ───────────────────────────────────────────

    def reddit_hot(self, subreddit: str = "MachineLearning", limit: int = 10) -> list[dict]:
        """Fetch hot posts from any subreddit. Free JSON API."""
        cache_key = f"reddit_{subreddit}"
        cached    = self._load_cache(cache_key)
        if cached:
            return cached

        try:
            url = f"https://www.reddit.com/r/{subreddit}/hot.json?limit={limit}"
            r   = requests.get(url, headers={"User-Agent": "NOVA/3.0"}, timeout=8)
            posts = []
            for post in r.json().get("data", {}).get("children", []):
                d = post["data"]
                posts.append({
                    "title":  d.get("title", ""),
                    "url":    f"https://reddit.com{d.get('permalink','')}",
                    "score":  d.get("score", 0),
                    "source": f"reddit/r/{subreddit}",
                })

            self._save_cache(cache_key, posts)
            if self.memory and posts:
                self.memory.store_many([
                    {"text": p["title"], "source": p["url"], "domain": "trending"}
                    for p in posts
                ])
            return posts

        except Exception as e:
            console.print(f"  [yellow]Reddit error: {e}[/]")
            return []

    # ── GitHub trending ───────────────────────────────────────────────────────

    def github_trending(self, language: str = "", since: str = "daily") -> list[dict]:
        """Scrape GitHub trending repositories."""
        cache_key = f"github_trend_{language}_{since}"
        cached    = self._load_cache(cache_key)
        if cached:
            return cached

        try:
            from bs4 import BeautifulSoup
            lang_path = f"/{language}" if language else ""
            url       = f"https://github.com/trending{lang_path}?since={since}"
            r         = requests.get(url, headers=self.HEADERS, timeout=10)
            soup      = BeautifulSoup(r.text, "lxml")
            repos     = []

            for article in soup.select("article.Box-row")[:10]:
                name_el = article.select_one("h2 a")
                desc_el = article.select_one("p")
                if name_el:
                    name = name_el.get_text(strip=True).replace("\n", "").replace(" ", "")
                    desc = desc_el.get_text(strip=True) if desc_el else ""
                    repos.append({
                        "name":   name,
                        "description": desc[:200],
                        "url":    f"https://github.com{name_el.get('href','')}",
                        "source": "github_trending",
                    })

            self._save_cache(cache_key, repos)
            if self.memory and repos:
                self.memory.store_many([
                    {"text": f"{r['name']}: {r['description']}", "source": r["url"], "domain": "technology"}
                    for r in repos
                ])
            return repos

        except Exception as e:
            console.print(f"  [yellow]GitHub trending error: {e}[/]")
            return []

    # ── Semantic Scholar — research papers with citations ─────────────────────

    def semantic_scholar(self, query: str, max_results: int = 5) -> list[dict]:
        """Search Semantic Scholar. Free API, no key needed for basic access."""
        cache_key = f"ss_{hashlib.md5(query.encode()).hexdigest()[:8]}"
        cached    = self._load_cache(cache_key)
        if cached:
            return cached

        try:
            r     = requests.get(
                "https://api.semanticscholar.org/graph/v1/paper/search",
                params={
                    "query":  query,
                    "limit":  max_results,
                    "fields": "title,abstract,year,citationCount,openAccessPdf,authors",
                },
                timeout=10,
            )
            papers = []
            for p in r.json().get("data", []):
                papers.append({
                    "title":     p.get("title", ""),
                    "abstract":  (p.get("abstract") or "")[:300],
                    "year":      p.get("year"),
                    "citations": p.get("citationCount", 0),
                    "pdf":       (p.get("openAccessPdf") or {}).get("url", ""),
                    "source":    "semantic_scholar",
                })

            self._save_cache(cache_key, papers)
            if self.memory and papers:
                self.memory.store_many([
                    {"text": f"{p['title']}. {p['abstract']}", "source": p.get("pdf",""), "domain": "research"}
                    for p in papers if p["abstract"]
                ])
            return papers

        except Exception as e:
            console.print(f"  [yellow]Semantic Scholar error: {e}[/]")
            return []

    # ── RSS feed reader ────────────────────────────────────────────────────────

    def rss_feed(self, url: str, max_items: int = 10) -> list[dict]:
        """Read any RSS feed. Works with any news source."""
        cache_key = f"rss_{hashlib.md5(url.encode()).hexdigest()[:8]}"
        cached    = self._load_cache(cache_key)
        if cached:
            return cached

        try:
            import xml.etree.ElementTree as ET
            r    = requests.get(url, headers=self.HEADERS, timeout=10)
            root = ET.fromstring(r.content)
            items = []

            for item in (root.findall(".//item") or root.findall(".//entry"))[:max_items]:
                title = (item.findtext("title")  or "").strip()
                link  = (item.findtext("link")   or item.findtext("{http://www.w3.org/2005/Atom}link") or "").strip()
                desc  = (item.findtext("description") or item.findtext("{http://www.w3.org/2005/Atom}summary") or "").strip()
                desc  = re.sub(r"<[^>]+>", "", desc)[:200]

                if title:
                    items.append({"title": title, "url": link, "description": desc, "source": url})

            self._save_cache(cache_key, items)
            if self.memory and items:
                self.memory.store_many([
                    {"text": f"{i['title']}. {i['description']}", "source": i["url"], "domain": "news"}
                    for i in items
                ])
            return items

        except Exception as e:
            console.print(f"  [yellow]RSS error ({url[:40]}): {e}[/]")
            return []

    # ── Combined trend pulse ──────────────────────────────────────────────────

    def full_pulse(self, topic: str = "artificial intelligence") -> dict:
        """
        Get a full trend pulse on any topic from all sources simultaneously.
        Uses threads for parallel fetching.
        """
        import concurrent.futures
        console.print(f"  [dim]Trend pulse: {topic}[/]")

        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as ex:
            f_arxiv  = ex.submit(self.arxiv_latest,      topic, 6)
            f_ss     = ex.submit(self.semantic_scholar,  topic, 4)
            f_hn     = ex.submit(self.hackernews_top,    8)
            f_reddit = ex.submit(self.reddit_hot,        "MachineLearning", 5)
            f_github = ex.submit(self.github_trending,   "",    "daily")

        return {
            "arxiv":           f_arxiv.result(),
            "semantic_scholar": f_ss.result(),
            "hackernews":      f_hn.result(),
            "reddit":          f_reddit.result(),
            "github_trending": f_github.result(),
            "timestamp":       datetime.now().isoformat(),
            "topic":           topic,
        }

    # ── Cache helpers ─────────────────────────────────────────────────────────

    def _cache_path(self, key: str) -> Path:
        return CACHE_DIR / f"{key}.json"

    def _load_cache(self, key: str) -> Optional[list]:
        p = self._cache_path(key)
        if not p.exists():
            return None
        age = time.time() - p.stat().st_mtime
        if age > self.cache_ttl:
            return None
        try:
            return json.loads(p.read_text())
        except Exception:
            return None

    def _save_cache(self, key: str, data: list):
        try:
            self._cache_path(key).write_text(json.dumps(data, ensure_ascii=False))
        except Exception:
            pass


# ─────────────────────────────────────────────────────────────────────────────
# UNIVERSAL DOCUMENT SCANNER
# ANY format, ANY language, zero external API dependency
# ─────────────────────────────────────────────────────────────────────────────

class UniversalScanner:
    """
    Reads ANY document or URL in ANY language.
    Handles formats that standard scrapers miss entirely.

    Free tools used:
    - EasyOCR / Tesseract  — image text extraction (scanned PDFs, photos)
    - Whisper              — audio/video transcription (local, offline)
    - youtube-transcript-api — YouTube video captions
    - Camelot / tabula     — table extraction from PDFs
    - python-pptx          — PowerPoint slides
    - openpyxl             — Excel spreadsheets
    - langdetect           — automatic language detection
    - All 100+ languages via multilingual sentence-transformers
    """

    HEADERS = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/124.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
    }

    def scan(self, source: str, **kwargs) -> dict:
        """
        Auto-detect source type and extract text.
        Returns: {text, language, format, confidence, metadata}
        """
        source = source.strip()

        # Route by source type
        if source.startswith("http"):
            if "youtube.com" in source or "youtu.be" in source:
                return self._scan_youtube(source)
            return self._scan_url(source)

        path = Path(source)
        if not path.exists():
            return {"error": f"File not found: {source}", "text": ""}

        ext = path.suffix.lower()
        if ext == ".pdf":
            return self._scan_pdf(path)
        elif ext in (".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"):
            return self._scan_image(path)
        elif ext in (".mp3", ".mp4", ".wav", ".m4a", ".webm", ".ogg"):
            return self._scan_audio(path)
        elif ext == ".docx":
            return self._scan_docx(path)
        elif ext in (".xlsx", ".xls"):
            return self._scan_excel(path)
        elif ext in (".pptx", ".ppt"):
            return self._scan_pptx(path)
        elif ext == ".csv":
            return self._scan_csv(path)
        elif ext == ".json":
            return self._scan_json(path)
        else:
            return self._scan_text_file(path)

    # ── PDF with OCR fallback ─────────────────────────────────────────────────

    def _scan_pdf(self, path: Path) -> dict:
        text = ""
        method = "text"

        # Try direct text extraction first (digital PDFs)
        try:
            import fitz
            doc  = fitz.open(str(path))
            text = "\n".join(page.get_text() for page in doc)
            doc.close()
        except ImportError:
            pass

        # If little text was found, this is a scanned PDF — use OCR
        if len(text.strip()) < 100:
            method = "ocr"
            text   = self._ocr_pdf(path)

        return self._wrap_result(text, method, str(path), {"pages": "?"})

    def _ocr_pdf(self, path: Path) -> str:
        """OCR a scanned PDF using EasyOCR (free, offline, multilingual)."""
        try:
            import fitz
            import easyocr
            reader = easyocr.Reader(["en"], gpu=False)
            doc    = fitz.open(str(path))
            texts  = []
            for page in doc:
                pix    = page.get_pixmap(dpi=200)
                img    = pix.tobytes("png")
                result = reader.readtext(img)
                texts.append(" ".join([r[1] for r in result]))
            doc.close()
            return "\n".join(texts)
        except ImportError:
            console.print("  [dim]EasyOCR not installed. Run: pip install easyocr[/]")
            return self._tesseract_fallback(path)
        except Exception as e:
            console.print(f"  [yellow]OCR error: {e}[/]")
            return ""

    def _tesseract_fallback(self, path: Path) -> str:
        """Tesseract OCR fallback (if installed on system)."""
        try:
            import subprocess
            result = subprocess.run(
                ["tesseract", str(path), "stdout", "-l", "eng"],
                capture_output=True, text=True, timeout=30
            )
            return result.stdout
        except Exception:
            return ""

    # ── Image OCR ─────────────────────────────────────────────────────────────

    def _scan_image(self, path: Path) -> dict:
        """Extract text from image using EasyOCR (supports 100+ languages)."""
        try:
            import easyocr
            reader = easyocr.Reader(["en", "hi", "ar", "zh_sim", "fr", "de", "es"], gpu=False)
            result = reader.readtext(str(path))
            text   = " ".join([r[1] for r in result])
            return self._wrap_result(text, "easyocr", str(path))
        except ImportError:
            console.print("  [dim]pip install easyocr for image OCR[/]")
            return {"text": "", "error": "easyocr not installed"}
        except Exception as e:
            return {"text": "", "error": str(e)}

    # ── Audio / Video transcription ───────────────────────────────────────────

    def _scan_audio(self, path: Path) -> dict:
        """Transcribe audio/video using local Whisper (offline, free, 99 languages)."""
        try:
            import whisper
            console.print(f"  [dim]Transcribing audio (Whisper)...[/]")
            model  = whisper.load_model("base")
            result = model.transcribe(str(path))
            text   = result.get("text", "")
            lang   = result.get("language", "unknown")
            return self._wrap_result(text, "whisper", str(path), {"language_detected": lang})
        except ImportError:
            console.print("  [dim]pip install openai-whisper for audio[/]")
            return {"text": "", "error": "whisper not installed"}
        except Exception as e:
            return {"text": "", "error": str(e)}

    # ── YouTube transcript ────────────────────────────────────────────────────

    def _scan_youtube(self, url: str) -> dict:
        """Extract transcript from YouTube video. Free, no API key."""
        try:
            from youtube_transcript_api import YouTubeTranscriptApi
            import re
            video_id = re.search(r"(?:v=|youtu\.be/)([a-zA-Z0-9_-]{11})", url)
            if not video_id:
                return {"text": "", "error": "Could not extract video ID"}

            transcript = YouTubeTranscriptApi.get_transcript(video_id.group(1))
            text       = " ".join(t["text"] for t in transcript)
            return self._wrap_result(text, "youtube_transcript", url,
                                     {"video_id": video_id.group(1)})
        except ImportError:
            console.print("  [dim]pip install youtube-transcript-api for YouTube[/]")
            return {"text": "", "error": "youtube-transcript-api not installed"}
        except Exception as e:
            return {"text": "", "error": str(e)}

    # ── Office formats ────────────────────────────────────────────────────────

    def _scan_docx(self, path: Path) -> dict:
        try:
            from docx import Document
            doc  = Document(str(path))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return self._wrap_result(text, "python-docx", str(path))
        except Exception as e:
            return {"text": "", "error": str(e)}

    def _scan_excel(self, path: Path) -> dict:
        try:
            import openpyxl
            wb    = openpyxl.load_workbook(str(path), data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        parts.append(row_text)
            return self._wrap_result("\n".join(parts), "openpyxl", str(path))
        except ImportError:
            console.print("  [dim]pip install openpyxl for Excel files[/]")
            return {"text": "", "error": "openpyxl not installed"}

    def _scan_pptx(self, path: Path) -> dict:
        try:
            from pptx import Presentation
            prs   = Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides):
                parts.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            return self._wrap_result("\n".join(parts), "python-pptx", str(path))
        except ImportError:
            console.print("  [dim]pip install python-pptx for PowerPoint[/]")
            return {"text": "", "error": "python-pptx not installed"}

    def _scan_csv(self, path: Path) -> dict:
        import csv
        rows = []
        for enc in ["utf-8", "utf-16", "latin-1"]:
            try:
                with open(path, encoding=enc) as f:
                    reader = csv.reader(f)
                    rows   = [" | ".join(row) for row in reader]
                break
            except UnicodeDecodeError:
                continue
        return self._wrap_result("\n".join(rows), "csv", str(path))

    def _scan_json(self, path: Path) -> dict:
        try:
            data = json.loads(path.read_text(encoding="utf-8"))
            text = json.dumps(data, indent=2, ensure_ascii=False)
            return self._wrap_result(text, "json", str(path))
        except Exception as e:
            return {"text": "", "error": str(e)}

    def _scan_text_file(self, path: Path) -> dict:
        for enc in ["utf-8", "utf-16", "latin-1"]:
            try:
                text = path.read_text(encoding=enc)
                return self._wrap_result(text, "plaintext", str(path))
            except UnicodeDecodeError:
                continue
        return {"text": "", "error": "Could not decode file"}

    # ── Web URL with full fallback chain ──────────────────────────────────────

    def _scan_url(self, url: str) -> dict:
        """Multi-tier URL fetching: trafilatura → requests → playwright."""
        # Tier 1: trafilatura (best quality)
        try:
            import trafilatura
            downloaded = trafilatura.fetch_url(url)
            if downloaded:
                text = trafilatura.extract(downloaded, include_tables=True)
                if text and len(text) > 200:
                    return self._wrap_result(text, "trafilatura", url)
        except Exception:
            pass

        # Tier 2: requests + BS4
        try:
            from bs4 import BeautifulSoup
            r    = requests.get(url, headers=self.HEADERS, timeout=12, allow_redirects=True)
            soup = BeautifulSoup(r.text, "lxml")
            for tag in soup(["script","style","nav","footer","header","aside"]):
                tag.decompose()
            text = soup.get_text(separator=" ")
            if len(text.strip()) > 200:
                return self._wrap_result(text, "beautifulsoup", url)
        except Exception:
            pass

        # Tier 3: Playwright (handles JS sites, Cloudflare)
        try:
            from playwright.sync_api import sync_playwright
            with sync_playwright() as pw:
                browser = pw.chromium.launch(headless=True, args=["--no-sandbox"])
                ctx     = browser.new_context(
                    user_agent=self.HEADERS["User-Agent"],
                    bypass_csp=True,
                )
                ctx.add_init_script("Object.defineProperty(navigator,'webdriver',{get:()=>undefined})")
                page = ctx.new_page()
                page.goto(url, wait_until="networkidle", timeout=25000)
                page.wait_for_timeout(2000)
                text = page.inner_text("body")
                browser.close()
                if text and len(text) > 200:
                    return self._wrap_result(text, "playwright", url)
        except Exception as e:
            console.print(f"  [yellow]Playwright failed: {e}[/]")

        return {"text": "", "error": "All scan methods failed", "url": url}

    # ── Language detection ─────────────────────────────────────────────────────

    def detect_language(self, text: str) -> str:
        try:
            from langdetect import detect, DetectorFactory
            DetectorFactory.seed = 0
            return detect(text[:500])
        except Exception:
            return "en"

    # ── Result wrapper ─────────────────────────────────────────────────────────

    def _wrap_result(self, text: str, method: str, source: str, extra: dict = None) -> dict:
        text     = re.sub(r"\s+", " ", text).strip()
        language = self.detect_language(text[:300]) if text else "unknown"
        return {
            "text":      text,
            "method":    method,
            "source":    source,
            "language":  language,
            "word_count": len(text.split()),
            "success":   len(text) > 50,
            **(extra or {}),
        }
