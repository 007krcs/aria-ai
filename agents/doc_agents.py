"""
ARIA — Document Agents
Four dedicated agents for handling documents end to end.

UploadAgent      → accepts files/URLs, validates, queues for processing
ProcessorAgent   → extracts text, cleans, chunks, detects language/topic
ReaderAgent      → answers questions FROM a specific document
KnowledgeAgent   → manages the knowledge base: search, stats, delete, summarise

Usage in main.py:
    from agents.doc_agents import UploadAgent, ProcessorAgent, ReaderAgent, KnowledgeAgent
"""

import os
import re
import time
import hashlib
import json
import shutil
from pathlib import Path
from datetime import datetime
from typing import Optional
from rich.console import Console
from rich.table import Table
from rich.progress import track

from core.engine import Engine
from core.memory import Memory
from tools.logger import Logger
from agents.agents import BaseAgent
from core.config import CHUNK_SIZE, CHUNK_OVERLAP, MODELS_DIR

console = Console()

# Where uploaded files are stored locally before processing
UPLOAD_DIR = Path(__file__).resolve().parent.parent / "data" / "uploads"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)


# ─────────────────────────────────────────────────────────────────────────────
# 1. UPLOAD AGENT
# ─────────────────────────────────────────────────────────────────────────────

class UploadAgent(BaseAgent):
    """
    Accepts any file or URL.
    Validates it, copies it to the uploads folder, returns a job dict
    that the ProcessorAgent can consume.

    Supported:
        Files  — .pdf  .docx  .txt  .md  .csv  .json  .html
        URLs   — any http/https URL
        Text   — raw string pasted directly
    """

    name = "upload_agent"

    SUPPORTED = {
        # Documents
        ".pdf", ".docx", ".txt", ".md", ".csv", ".json", ".html", ".rst",
        # Spreadsheets / Presentations
        ".xlsx", ".xls", ".pptx", ".ppt",
        # Audio
        ".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac", ".wma",
        # Video
        ".mp4", ".mkv", ".avi", ".mov", ".webm", ".flv", ".wmv",
        # Images (for OCR)
        ".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp",
    }

    AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac", ".wma"}
    VIDEO_EXTS = {".mp4", ".mkv", ".avi", ".mov", ".webm", ".flv", ".wmv"}
    IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"}

    def upload_file(self, path: str, domain: str = "general") -> dict:
        """
        Copy a local file into the uploads folder and return a job dict.
        """
        src = Path(path)
        if not src.exists():
            return self._err(f"File not found: {path}")
        if src.suffix.lower() not in self.SUPPORTED:
            return self._err(
                f"Unsupported format: {src.suffix}\n"
                f"Supported: documents, audio (mp3/wav/mp4/mkv...), images (jpg/png...)"
            )

        # Deduplicate by content hash
        content_hash = self._hash_file(src)
        dest          = UPLOAD_DIR / f"{content_hash}{src.suffix}"
        if not dest.exists():
            shutil.copy2(src, dest)
            console.print(f"  [green]Uploaded:[/] {src.name} → {dest}")
        else:
            console.print(f"  [dim]Already uploaded:[/] {src.name} (same content exists)")

        return {
            "status":       "ready",
            "type":         "file",
            "source":       str(dest),
            "original_name": src.name,
            "domain":       domain,
            "content_hash": content_hash,
            "queued_at":    datetime.now().isoformat(),
        }

    def upload_url(self, url: str, domain: str = "general") -> dict:
        """Register a URL for processing."""
        if not url.startswith(("http://", "https://")):
            return self._err(f"Invalid URL: {url}")
        console.print(f"  [green]URL queued:[/] {url[:70]}")
        return {
            "status":    "ready",
            "type":      "url",
            "source":    url,
            "domain":    domain,
            "queued_at": datetime.now().isoformat(),
        }

    def upload_text(self, text: str, title: str = "paste", domain: str = "general") -> dict:
        """Accept raw pasted text."""
        if len(text.strip()) < 20:
            return self._err("Text too short (minimum 20 characters)")
        # Save to file so processor can handle it uniformly
        content_hash = hashlib.md5(text.encode()).hexdigest()[:12]
        dest         = UPLOAD_DIR / f"{content_hash}_{title}.txt"
        dest.write_text(text, encoding="utf-8")
        console.print(f"  [green]Text saved:[/] {len(text)} chars → {dest.name}")
        return {
            "status":    "ready",
            "type":      "text",
            "source":    str(dest),
            "domain":    domain,
            "queued_at": datetime.now().isoformat(),
        }

    def list_uploads(self) -> list[dict]:
        """List all files currently in the uploads folder."""
        files = []
        for f in UPLOAD_DIR.iterdir():
            if f.name.startswith("."):
                continue
            files.append({
                "name":     f.name,
                "size_kb":  round(f.stat().st_size / 1024, 1),
                "modified": datetime.fromtimestamp(f.stat().st_mtime).strftime("%Y-%m-%d %H:%M"),
            })
        return sorted(files, key=lambda x: x["modified"], reverse=True)

    def _hash_file(self, path: Path) -> str:
        h = hashlib.md5()
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()[:16]

    def _err(self, msg: str) -> dict:
        console.print(f"  [red]Upload error:[/] {msg}")
        return {"status": "error", "error": msg}


# ─────────────────────────────────────────────────────────────────────────────
# 2. PROCESSOR AGENT
# ─────────────────────────────────────────────────────────────────────────────

class ProcessorAgent(BaseAgent):
    """
    Takes an upload job dict (from UploadAgent) and:
      1. Extracts raw text from the file/URL
      2. Detects language and auto-assigns domain/topic using the LLM
      3. Cleans and chunks the text
      4. Stores all chunks in ChromaDB memory
      5. Generates a summary and stores it too
      6. Returns a processing report

    This is the heart of the knowledge base pipeline.
    """

    name = "processor_agent"

    def process(self, job: dict) -> dict:
        """
        Process one upload job. Returns a report dict.
        job = output of UploadAgent.upload_file / upload_url / upload_text
        """
        if job.get("status") == "error":
            return job

        source = job["source"]
        domain = job.get("domain", "general")

        console.print(f"\n[bold]Processor:[/] {source[:60]}")

        # ── Step 1: Extract raw text ─────────────────────────────────────────
        console.print("  [dim]Step 1/5 — Extracting text...[/]")
        raw_text = self._extract(source, job["type"])
        if not raw_text or len(raw_text.strip()) < 50:
            return {"status": "error", "error": "Could not extract readable text"}

        word_count = len(raw_text.split())
        console.print(f"  [dim]Extracted {word_count} words[/]")

        # ── Step 2: Detect language ──────────────────────────────────────────
        console.print("  [dim]Step 2/5 — Detecting language...[/]")
        language = self._detect_language(raw_text[:500])

        # ── Step 3: Auto-detect topic/domain with LLM ───────────────────────
        console.print("  [dim]Step 3/5 — Classifying topic...[/]")
        if domain == "general":
            domain = self._detect_domain(raw_text[:600])
        console.print(f"  [dim]Domain: {domain} | Language: {language}[/]")

        # ── Step 4: Clean, chunk, store ──────────────────────────────────────
        console.print("  [dim]Step 4/5 — Chunking and storing in memory...[/]")
        cleaned = self._clean(raw_text)
        chunks  = self._chunk(cleaned)

        chunk_dicts = [
            {
                "text":   c,
                "source": source,
                "domain": domain,
                "metadata": {
                    "language": language,
                    "original_name": job.get("original_name", source),
                    "processed_at":  datetime.now().isoformat(),
                }
            }
            for c in chunks
        ]
        self.memory.store_many(chunk_dicts)
        self.logger.log_ingestion(source, job["type"], len(chunks), domain)

        # ── Step 5: Generate and store summary ───────────────────────────────
        console.print("  [dim]Step 5/5 — Generating document summary...[/]")
        summary = self._summarise(raw_text[:2000], domain, language)

        # Store summary as a special high-priority chunk
        self.memory.store(
            text=f"[SUMMARY of {job.get('original_name', source)}] {summary}",
            source=source,
            domain=domain,
            metadata={"type": "summary", "language": language},
        )

        report = {
            "status":      "done",
            "source":      source,
            "domain":      domain,
            "language":    language,
            "word_count":  word_count,
            "chunks":      len(chunks),
            "summary":     summary,
            "processed_at": datetime.now().isoformat(),
        }
        console.print(f"  [green]Done:[/] {len(chunks)} chunks stored | domain={domain}")
        return report

    def process_batch(self, jobs: list[dict]) -> list[dict]:
        """Process multiple upload jobs in sequence."""
        reports = []
        for i, job in enumerate(jobs, 1):
            console.print(f"\n[dim]Batch {i}/{len(jobs)}[/]")
            reports.append(self.process(job))
        return reports

    # ── Extraction ────────────────────────────────────────────────────────────

    def _extract(self, source: str, source_type: str) -> str:
        if source_type == "url":
            return self._extract_url(source)
        else:
            return self._extract_file(source)

    def _extract_file(self, path: str) -> str:
        suffix = Path(path).suffix.lower()
        try:
            if suffix == ".pdf":
                import fitz
                doc        = fitz.open(path)
                full_text  = "\n".join(p.get_text() for p in doc)
                doc.close()
                # If little text extracted, it's a scanned PDF — use TieredOCR
                if len(full_text.strip()) < 100:
                    console.print("  [dim]Scanned PDF detected — using TieredOCR[/]")
                    try:
                        from agents.tiered_ocr import TieredOCR
                        ocr_result = TieredOCR().read_pdf(path)
                        return ocr_result.get("text", full_text)
                    except Exception as ocr_err:
                        console.print(f"  [yellow]TieredOCR fallback error: {ocr_err}[/]")
                return full_text
            elif suffix == ".docx":
                from docx import Document
                doc = Document(path)
                return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            elif suffix in (".csv",):
                import csv
                rows = []
                with open(path, encoding="utf-8", errors="replace") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        rows.append(" | ".join(row))
                return "\n".join(rows)
            elif suffix == ".json":
                with open(path, encoding="utf-8") as f:
                    data = json.load(f)
                return json.dumps(data, indent=2, ensure_ascii=False)
            elif suffix == ".html":
                from bs4 import BeautifulSoup
                html = Path(path).read_text(encoding="utf-8", errors="replace")
                soup = BeautifulSoup(html, "lxml")
                for tag in soup(["script","style","nav","footer"]):
                    tag.decompose()
                return soup.get_text(separator=" ")
            elif suffix in self.AUDIO_EXTS or suffix in self.VIDEO_EXTS:
                return self._extract_audio_video(path)
            elif suffix in self.IMAGE_EXTS:
                return self._extract_image_ocr(path)
            elif suffix in (".xlsx", ".xls"):
                try:
                    import openpyxl
                    wb    = openpyxl.load_workbook(path, data_only=True)
                    parts = []
                    for sheet in wb.worksheets:
                        parts.append(f"Sheet: {sheet.title}")
                        for row in sheet.iter_rows(values_only=True):
                            row_text = " | ".join(str(c) for c in row if c is not None)
                            if row_text.strip():
                                parts.append(row_text)
                    return "\n".join(parts)
                except ImportError:
                    console.print("  [dim]pip install openpyxl for Excel files[/]")
            elif suffix in (".pptx", ".ppt"):
                try:
                    from pptx import Presentation
                    prs   = Presentation(path)
                    parts = []
                    for i, slide in enumerate(prs.slides):
                        parts.append(f"--- Slide {i+1} ---")
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                parts.append(shape.text.strip())
                    return "\n".join(parts)
                except ImportError:
                    console.print("  [dim]pip install python-pptx for PowerPoint files[/]")
            else:  # .txt .md .rst
                for enc in ["utf-8", "utf-16", "latin-1"]:
                    try:
                        return Path(path).read_text(encoding=enc)
                    except UnicodeDecodeError:
                        continue
        except ImportError as e:
            console.print(f"  [yellow]Missing package:[/] {e}")
        except Exception as e:
            console.print(f"  [red]Extraction error:[/] {e}")
        return ""

    def _extract_audio_video(self, path: str) -> str:
        """
        Transcribe audio or video file using local Whisper.
        Works completely offline. Supports 99 languages.

        Models (trade-off speed vs accuracy):
          tiny   — fastest, ~1GB RAM, good for English
          base   — fast,    ~1GB RAM, good multilingual  ← default
          small  — slower,  ~2GB RAM, better accuracy
          medium — slow,    ~5GB RAM, near-perfect
          large  — slowest, ~10GB RAM, best quality

        Install: pip install openai-whisper
        Also needs ffmpeg: https://ffmpeg.org/download.html
        Windows: winget install ffmpeg
        """
        import time
        t0 = time.time()
        console.print(f"  [dim]Transcribing: {Path(path).name}[/]")
        console.print(f"  [dim]This may take a few minutes for long files...[/]")

        # Method 1: openai-whisper (best quality, runs fully offline)
        try:
            import whisper
            model_size = "base"   # change to "small" or "medium" for better quality
            console.print(f"  [dim]Loading Whisper {model_size} model...[/]")
            model  = whisper.load_model(model_size)
            result = model.transcribe(path, verbose=False)
            text   = result.get("text", "").strip()
            lang   = result.get("language", "unknown")
            ms     = int((time.time() - t0) * 1000)
            console.print(
                f"  [green]Whisper transcribed:[/] {len(text.split())} words "
                f"in {lang} ({ms//1000}s)"
            )
            return text
        except ImportError:
            console.print(
                "  [yellow]openai-whisper not installed.[/]\n"
                "  Run: pip install openai-whisper\n"
                "  Also needs ffmpeg: winget install ffmpeg (Windows)"
            )

        # Method 2: faster-whisper (faster, less RAM)
        try:
            from faster_whisper import WhisperModel
            console.print("  [dim]Using faster-whisper...[/]")
            model  = WhisperModel("base", device="cpu", compute_type="int8")
            segs, info = model.transcribe(path)
            text   = " ".join(seg.text for seg in segs).strip()
            ms     = int((time.time() - t0) * 1000)
            console.print(
                f"  [green]faster-whisper:[/] {len(text.split())} words "
                f"in {info.language} ({ms//1000}s)"
            )
            return text
        except ImportError:
            pass

        # Method 3: Ollama speech (if a speech model is installed)
        # Currently experimental — most Ollama models don't support audio yet
        console.print(
            "  [red]No transcription engine found.[/]\n"
            "  Install one of:\n"
            "    pip install openai-whisper   (+ ffmpeg)\n"
            "    pip install faster-whisper"
        )
        return ""

    def _extract_image_ocr(self, path: str) -> str:
        """Extract text from image using TieredOCR (6-tier fallback)."""
        try:
            from agents.tiered_ocr import TieredOCR
            result = TieredOCR().read(path)
            return result.text
        except Exception as e:
            console.print(f"  [yellow]Image OCR error: {e}[/]")
            # Fallback to vision LLM
            try:
                from agents.vision_ocr import VisionOCR
                import base64
                b64 = base64.b64encode(Path(path).read_bytes()).decode()
                result = VisionOCR().image_to_text(b64)
                return result.get("text", "")
            except Exception:
                return ""

    BROWSER_HEADERS = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/124.0.0.0 Safari/537.36"
        ),
        "Accept":          "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.9",
        "Accept-Encoding": "gzip, deflate, br",
        "Connection":      "keep-alive",
        "Upgrade-Insecure-Requests": "1",
        "Sec-Fetch-Dest":  "document",
        "Sec-Fetch-Mode":  "navigate",
        "Sec-Fetch-Site":  "none",
    }

    def _extract_url(self, url: str) -> str:
        # Tier 1: trafilatura (best text extraction quality)
        try:
            import trafilatura
            downloaded = trafilatura.fetch_url(url)
            if downloaded:
                text = trafilatura.extract(downloaded, include_tables=True, include_links=False)
                if text and len(text) > 200:
                    return text
        except Exception:
            pass

        # Tier 2: requests with browser headers
        try:
            import requests
            from bs4 import BeautifulSoup
            r = requests.get(url, timeout=12, headers=self.BROWSER_HEADERS, allow_redirects=True)
            if r.status_code == 200 and len(r.text) > 200:
                soup = BeautifulSoup(r.text, "lxml")
                for tag in soup(["script","style","nav","footer","header"]):
                    tag.decompose()
                text = soup.get_text(separator=" ")
                if len(text.strip()) > 200:
                    return text
        except Exception:
            pass

        # Tier 3: requests session with cookie warm-up
        try:
            import requests
            from urllib.parse import urlparse
            from bs4 import BeautifulSoup
            session = requests.Session()
            session.headers.update(self.BROWSER_HEADERS)
            base = f"{urlparse(url).scheme}://{urlparse(url).netloc}"
            session.get(base, timeout=8)
            r2 = session.get(url, timeout=12, allow_redirects=True)
            if r2.status_code == 200:
                soup = BeautifulSoup(r2.text, "lxml")
                for tag in soup(["script","style","nav","footer","header"]):
                    tag.decompose()
                return soup.get_text(separator=" ")
        except Exception:
            pass

        # Tier 4: Playwright headless browser (handles Cloudflare, JS SPAs)
        console.print(f"  [dim]Trying Playwright for {url[:60]}...[/]")
        html = self._playwright_fetch(url)
        if html:
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html, "lxml")
                for tag in soup(["script","style","nav","footer","header"]):
                    tag.decompose()
                return soup.get_text(separator=" ")
            except Exception:
                return html

        console.print(f"  [red]All fetch methods failed for:[/] {url}")
        return ""

    def _playwright_fetch(self, url: str) -> str | None:
        """Full stealth Playwright fetch — passes Cloudflare and JS-heavy sites."""
        try:
            from playwright.sync_api import sync_playwright, TimeoutError as PWTimeout
        except ImportError:
            console.print("  [dim]pip install playwright && playwright install chromium[/]")
            return None

        STEALTH = """
            Object.defineProperty(navigator,'webdriver',{get:()=>undefined});
            window.chrome={runtime:{connect:()=>{},sendMessage:()=>{},onMessage:{addListener:()=>{}}}};
            Object.defineProperty(navigator,'plugins',{get:()=>{const a=[{name:'Chrome PDF Plugin'},{name:'Chrome PDF Viewer'}];a.__proto__=PluginArray.prototype;return a;}});
            Object.defineProperty(navigator,'languages',{get:()=>['en-US','en']});
            Object.defineProperty(navigator,'hardwareConcurrency',{get:()=>8});
            Object.defineProperty(navigator,'platform',{get:()=>'Win32'});
            const gp=WebGLRenderingContext.prototype.getParameter;
            WebGLRenderingContext.prototype.getParameter=function(p){if(p===37445)return 'Intel Inc.';if(p===37446)return 'Intel Iris OpenGL Engine';return gp.call(this,p);};
        """
        from urllib.parse import urlparse as _up
        try:
            with sync_playwright() as pw:
                browser = pw.chromium.launch(
                    headless=True,
                    args=["--no-sandbox","--disable-blink-features=AutomationControlled",
                          "--disable-dev-shm-usage","--window-size=1280,800"],
                )
                ctx = browser.new_context(
                    viewport={"width":1280,"height":800},
                    user_agent=self.BROWSER_HEADERS["User-Agent"],
                    locale="en-US", timezone_id="America/Chicago",
                    bypass_csp=True,
                    extra_http_headers={
                        "Accept":"text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
                        "Accept-Language":"en-US,en;q=0.9",
                        "Sec-Fetch-Dest":"document","Sec-Fetch-Mode":"navigate",
                        "Sec-Fetch-Site":"none","Sec-Fetch-User":"?1",
                    },
                )
                ctx.add_init_script(STEALTH)
                page = ctx.new_page()
                page.route("**/*.{png,jpg,jpeg,gif,svg,ico,woff,woff2,ttf}",
                           lambda r: r.abort())
                try:
                    # Visit homepage first for cookie/trust warm-up
                    parsed = _up(url)
                    base = f"{parsed.scheme}://{parsed.netloc}"
                    if url.rstrip("/") != base:
                        try:
                            page.goto(base, wait_until="domcontentloaded", timeout=12000)
                            page.wait_for_timeout(1200)
                        except Exception:
                            pass

                    try:
                        page.goto(url, wait_until="networkidle", timeout=28000)
                    except PWTimeout:
                        page.goto(url, wait_until="domcontentloaded", timeout=18000)

                    try:
                        page.wait_for_selector("main, article, p, h1", timeout=8000)
                    except Exception:
                        page.wait_for_timeout(2500)

                    page.evaluate("window.scrollTo(0, document.body.scrollHeight*0.5)")
                    page.wait_for_timeout(1200)

                    html = page.content()
                    console.print(f"  [green]Playwright:[/] {len(html):,} chars")
                    return html if len(html) > 200 else None
                except PWTimeout:
                    try: return page.content()
                    except: return None
                finally:
                    browser.close()
        except Exception as e:
            console.print(f"  [yellow]Playwright error: {e}[/]")
            return None

    # ── Cleaning ──────────────────────────────────────────────────────────────

    def _clean(self, text: str) -> str:
        text = re.sub(r"\s+", " ", text)
        text = re.sub(r"[^\x20-\x7E\u0900-\u097F\u0600-\u06FF\u4e00-\u9fff]", " ", text)
        return text.strip()

    # ── Chunking ──────────────────────────────────────────────────────────────

    def _chunk(self, text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
        words  = text.split()
        chunks = []
        i      = 0
        while i < len(words):
            chunk = " ".join(words[i: i + size])
            if chunk.strip():
                chunks.append(chunk)
            i += size - overlap
        return chunks

    # ── Language detection ────────────────────────────────────────────────────

    def _detect_language(self, text: str) -> str:
        try:
            from langdetect import detect
            code = detect(text)
            return {
                "hi": "Hindi", "en": "English", "fr": "French",
                "de": "German", "es": "Spanish", "ar": "Arabic",
                "zh-cn": "Chinese", "ja": "Japanese", "bn": "Bengali",
                "ta": "Tamil",  "te": "Telugu",  "ur": "Urdu",
            }.get(code, code)
        except Exception:
            return "English"

    # ── LLM domain detection ──────────────────────────────────────────────────

    def _detect_domain(self, text: str) -> str:
        prompt = (
            f"Read this text and classify it into ONE domain from this list:\n"
            f"technology, science, medicine, law, finance, history, education, "
            f"news, sports, arts, religion, politics, engineering, business, general\n\n"
            f"Text: {text[:400]}\n\n"
            f"Reply with ONE word only — the domain name."
        )
        result = self.engine.generate(prompt, temperature=0.0)
        result = result.strip().lower().split()[0] if result.strip() else "general"
        valid  = {
            "technology","science","medicine","law","finance","history",
            "education","news","sports","arts","religion","politics",
            "engineering","business","general"
        }
        return result if result in valid else "general"

    # ── LLM summarisation ─────────────────────────────────────────────────────

    def _summarise(self, text: str, domain: str, language: str) -> str:
        prompt = (
            f"Summarise this {domain} document in 3-4 sentences. "
            f"Be specific — mention key topics, names, and conclusions.\n\n"
            f"Document:\n{text[:1500]}\n\n"
            f"Summary:"
        )
        return self.engine.generate(prompt, temperature=0.2)


# ─────────────────────────────────────────────────────────────────────────────
# 3. READER AGENT
# ─────────────────────────────────────────────────────────────────────────────

class ReaderAgent(BaseAgent):
    """
    Reads and answers questions from a SPECIFIC document.
    Unlike the general Researcher (which searches all memory),
    the Reader focuses only on one document's chunks.

    Use this when the user says:
        "In the PDF I uploaded, what does it say about X?"
        "From the report, what are the key findings?"
        "Read the contract and find the termination clause."
    """

    name = "reader_agent"

    def answer(self, question: str, source: str) -> dict:
        """
        Answer a question using only chunks from a specific source document.

        Args:
            question: what the user wants to know
            source:   the source string (file path or URL) stored in memory

        Returns:
            { answer, confidence, relevant_chunks, source }
        """
        console.print(f"\n[bold]Reader:[/] searching '{Path(source).name[:40]}'")

        # Get all chunks from this specific source
        all_results = self.memory.search(question, top_k=10, min_similarity=0.0)
        doc_chunks  = [r for r in all_results if r["source"] == source]

        if not doc_chunks:
            # Try searching by partial source match (URL may have changed)
            doc_chunks = [r for r in all_results if source[:30] in r["source"]]

        if not doc_chunks:
            return {
                "answer":           "Document not found in memory. Please ingest it first.",
                "confidence":        0.0,
                "relevant_chunks":  0,
                "source":           source,
            }

        # Take top 4 most relevant chunks from this document
        top_chunks = sorted(doc_chunks, key=lambda x: x["similarity"], reverse=True)[:4]
        context    = "\n\n".join(f"[Chunk {i+1}]\n{c['text']}" for i, c in enumerate(top_chunks))

        console.print(f"  [dim]Found {len(doc_chunks)} chunks, using top {len(top_chunks)}[/]")

        prompt = (
            f"You are reading a specific document. Answer the question using ONLY "
            f"the document excerpts below. If the answer is not in the excerpts, say so.\n\n"
            f"Document excerpts:\n{context}\n\n"
            f"Question: {question}\n\n"
            f"Answer:"
        )

        answer     = self.engine.generate(prompt, temperature=0.1)
        confidence = top_chunks[0]["similarity"] if top_chunks else 0.0

        self.logger.log_agent_run(
            agent_name=self.name,
            task=question,
            result=answer[:200],
            score=confidence,
        )

        return {
            "answer":          answer,
            "confidence":      round(confidence, 3),
            "relevant_chunks": len(doc_chunks),
            "source":          source,
        }

    def extract_sections(self, source: str, sections: list[str]) -> dict:
        """
        Extract specific sections from a document.
        e.g. sections=["introduction", "conclusion", "methodology"]
        """
        results = {}
        for section in sections:
            result        = self.answer(f"What does the document say about {section}?", source)
            results[section] = result["answer"]
        return results

    def list_document_topics(self, source: str) -> str:
        """Ask the LLM what topics this document covers."""
        chunks = self.memory.search("main topics summary", top_k=5, min_similarity=0.0)
        doc_chunks = [c for c in chunks if source[:30] in c["source"]]
        if not doc_chunks:
            return "Document not found in memory."
        context = "\n".join(c["text"][:200] for c in doc_chunks[:3])
        prompt  = (
            f"Based on these excerpts, list the main topics this document covers "
            f"(bullet points, max 8 topics):\n\n{context}\n\nTopics:"
        )
        return self.engine.generate(prompt, temperature=0.2)


# ─────────────────────────────────────────────────────────────────────────────
# 4. KNOWLEDGE BASE AGENT
# ─────────────────────────────────────────────────────────────────────────────

class KnowledgeAgent(BaseAgent):
    """
    Manages the entire knowledge base.
    Provides search, stats, summaries, deduplication, and cleanup.

    Think of this as the librarian — it knows what's in the library,
    can find anything, and keeps the collection organised.
    """

    name = "knowledge_agent"

    def search(self, query: str, domain: Optional[str] = None, top_k: int = 5) -> list[dict]:
        """Search the knowledge base and return ranked results."""
        return self.memory.search(query, top_k=top_k, domain=domain)

    def ask(self, question: str, domain: Optional[str] = None) -> dict:
        """
        Answer any question from the entire knowledge base.
        Smarter than a raw search — uses LLM to synthesise across chunks.
        """
        console.print(f"\n[bold]KnowledgeAgent:[/] {question[:60]}")

        hits = self.memory.search(question, top_k=5, domain=domain)
        if not hits:
            return {"answer": "No relevant knowledge found. Ingest some documents first.", "sources": []}

        context = "\n\n".join(
            f"[{h['source'].split('/')[-1][:30]} | relevance {h['similarity']}]\n{h['text']}"
            for h in hits
        )
        prompt = (
            f"Answer the question using the knowledge base excerpts below.\n"
            f"Synthesise across sources. Be specific.\n\n"
            f"Knowledge:\n{context[:1800]}\n\n"
            f"Question: {question}\n\nAnswer:"
        )
        answer  = self.engine.generate(prompt, temperature=0.2)
        sources = list({h["source"] for h in hits})

        return {"answer": answer, "sources": sources, "chunks_used": len(hits)}

    def stats(self) -> dict:
        """Full knowledge base statistics."""
        mem_stats  = self.memory.stats()
        log_stats  = self.logger.get_stats()
        return {**mem_stats, **log_stats}

    def print_stats(self):
        """Pretty-print knowledge base stats."""
        s     = self.stats()
        table = Table(title="Knowledge Base Stats", show_header=False)
        table.add_column("Metric", style="dim")
        table.add_column("Value",  style="bold")
        table.add_row("Total chunks stored",  str(s.get("total_chunks", 0)))
        table.add_row("Sources ingested",     str(s.get("ingested_sources", 0)))
        table.add_row("Total queries handled",str(s.get("total_interactions", 0)))
        table.add_row("Avg confidence",       str(s.get("avg_confidence", 0)))
        table.add_row("DB path",              s.get("db_path", ""))
        console.print(table)

    def list_sources(self) -> list[dict]:
        """List all ingested document sources from the logger."""
        with self.logger._connect() as conn:
            rows = conn.execute(
                "SELECT source, source_type, chunks, domain, ts FROM ingested_sources ORDER BY ts DESC"
            ).fetchall()
        return [dict(r) for r in rows]

    def print_sources(self):
        """Pretty-print all ingested sources."""
        sources = self.list_sources()
        if not sources:
            console.print("[yellow]No documents ingested yet.[/]")
            console.print("Run: python main.py ingest <file or URL>")
            return
        table = Table(title=f"Knowledge Base — {len(sources)} sources")
        table.add_column("#",       width=3)
        table.add_column("Source",  width=40)
        table.add_column("Type",    width=6)
        table.add_column("Domain",  width=12)
        table.add_column("Chunks",  width=7)
        table.add_column("Date",    width=16)
        for i, s in enumerate(sources, 1):
            name = s["source"].split("/")[-1].split("\\")[-1][:38]
            table.add_row(
                str(i), name, s["source_type"],
                s["domain"], str(s["chunks"]),
                s["ts"][:16],
            )
        console.print(table)

    def delete_source(self, source: str) -> bool:
        """Remove all chunks from a specific source."""
        self.memory.delete_by_source(source)
        return True

    def summarise_kb(self) -> str:
        """Ask the LLM to summarise everything in the knowledge base."""
        sources = self.list_sources()
        if not sources:
            return "Knowledge base is empty."
        source_list = "\n".join(
            f"- {s['source'].split('/')[-1]} ({s['domain']}, {s['chunks']} chunks)"
            for s in sources[:15]
        )
        prompt = (
            f"I have a knowledge base with these documents:\n{source_list}\n\n"
            f"Write a brief paragraph describing what topics this knowledge base covers "
            f"and what kinds of questions it can answer."
        )
        return self.engine.generate(prompt, temperature=0.3)

    def find_document(self, name_hint: str) -> list[dict]:
        """Find ingested sources matching a name hint."""
        sources = self.list_sources()
        return [s for s in sources if name_hint.lower() in s["source"].lower()]
